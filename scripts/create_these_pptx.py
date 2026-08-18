#!/usr/bin/env python3
"""
Soutenance de thèse — Agonistes du GLP-1 & maladie de Parkinson
~30 diapositives, fidèles au manuscrit (chiffres d’essais, tableaux 1–3, conclusions).
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0A, 0x36, 0x42)
TEAL = RGBColor(0x1A, 0x6B, 0x7A)
TEAL_LT = RGBColor(0x2A, 0x8A, 0x9A)
ACCENT = RGBColor(0xC4, 0x7A, 0x2C)
ACCENT_DK = RGBColor(0x9A, 0x5C, 0x1A)
LIGHT = RGBColor(0xF3, 0xF6, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x22, 0x26)
MUTED = RGBColor(0x4E, 0x5E, 0x66)
SOFT = RGBColor(0xE4, 0xEC, 0xEE)
ROW_ALT = RGBColor(0xEE, 0xF4, 0xF5)
WARN = RGBColor(0x8B, 0x45, 0x13)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 30
ASSETS = Path("/workspace/docs/presentation/assets")


def set_run(run, size=16, bold=False, color=DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)
    return shape


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def add_round(slide, left, top, width, height, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_textbox(slide, left, top, width, height, text, size=16, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color, font=font)
    return box


def add_bullets(slide, left, top, width, height, items, size=14, color=DARK, spacing=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = "▸  " + item[0]
            set_run(run, size=size, bold=True, color=color)
            run2 = p.add_run()
            run2.text = item[1]
            set_run(run2, size=size, bold=False, color=MUTED)
        else:
            run.text = "▸  " + item
            set_run(run, size=size, bold=False, color=color)
    return box


def add_footer(slide, page, total=TOTAL):
    add_rect(slide, 0, Inches(7.12), SLIDE_W, Inches(0.38), NAVY)
    add_textbox(
        slide, Inches(0.35), Inches(7.16), Inches(10.5), Inches(0.28),
        "Agonistes du GLP-1 — Mise au point et potentiel dans la maladie de Parkinson",
        size=10, color=WHITE,
    )
    add_textbox(
        slide, Inches(11.5), Inches(7.16), Inches(1.5), Inches(0.28),
        f"{page}  /  {total}", size=10, color=WHITE, align=PP_ALIGN.RIGHT,
    )


def header(slide, kicker, title, subtitle=None):
    add_bg(slide, LIGHT)
    h = Inches(1.15) if not subtitle else Inches(1.35)
    add_rect(slide, 0, 0, SLIDE_W, h, NAVY)
    add_rect(slide, 0, 0, Inches(0.12), h, ACCENT)
    add_textbox(slide, Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.28),
                kicker.upper(), size=11, bold=True, color=TEAL_LT)
    add_textbox(slide, Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.45),
                title, size=24, bold=True, color=WHITE, font="Georgia")
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.28),
                    subtitle, size=12, color=RGBColor(0xB0, 0xD0, 0xD8))


def card(slide, left, top, width, height, title, bullets, header_color=TEAL, body_size=13):
    add_round(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, width, Inches(0.42), header_color)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.06), width - Inches(0.3), Inches(0.32),
                title, size=13, bold=True, color=WHITE)
    add_bullets(
        slide, left + Inches(0.15), top + Inches(0.55),
        width - Inches(0.3), height - Inches(0.65),
        bullets, size=body_size, spacing=5,
    )


def kpi(slide, left, top, width, height, value, label, sub=None, bg=NAVY):
    add_round(slide, left, top, width, height, bg)
    add_textbox(
        slide, left + Inches(0.12), top + Inches(0.18), width - Inches(0.24), Inches(0.55),
        value, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font="Georgia",
    )
    add_textbox(
        slide, left + Inches(0.12), top + Inches(0.75), width - Inches(0.24), Inches(0.45),
        label, size=12, bold=True, color=RGBColor(0xB8, 0xD8, 0xE0), align=PP_ALIGN.CENTER,
    )
    if sub:
        add_textbox(
            slide, left + Inches(0.12), top + Inches(1.15), width - Inches(0.24), Inches(0.4),
            sub, size=11, color=RGBColor(0x90, 0xB8, 0xC0), align=PP_ALIGN.CENTER,
        )


def add_table(slide, left, top, width, rows, col_widths, font_size=11, row_h=0.38):
    n_rows, n_cols = len(rows), len(rows[0])
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, width, Inches(row_h * n_rows)
    )
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            is_header = r == 0
            set_run(run, size=font_size, bold=is_header, color=WHITE if is_header else DARK)
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = NAVY
            elif r % 2 == 0:
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def add_picture_fit(slide, path, left, top, width, height):
    pic = slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    return pic


def section_divider(prs, blank, num, title, subtitle, page):
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.2), SLIDE_H, ACCENT)
    add_textbox(s, Inches(1.0), Inches(2.2), Inches(11), Inches(0.4),
                f"PARTIE  {num}", size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(1.0), Inches(2.7), Inches(11), Inches(1.2),
                title, size=32, bold=True, color=WHITE, font="Georgia")
    add_textbox(s, Inches(1.0), Inches(4.2), Inches(11), Inches(0.8),
                subtitle, size=16, color=RGBColor(0xA8, 0xD0, 0xD8))
    add_textbox(s, Inches(1.0), Inches(6.5), Inches(11), Inches(0.3),
                f"{page}  /  {TOTAL}", size=11, color=RGBColor(0x70, 0x98, 0xA0))
    return s


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    p = 0

    # 1. TITRE
    p += 1
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
    add_rect(s, 0, Inches(5.85), SLIDE_W, Inches(1.65), TEAL)
    add_textbox(s, Inches(0.7), Inches(1.15), Inches(11.5), Inches(0.35),
                "SOUTENANCE DE THÈSE  ·  PHARMACIE", size=13, bold=True, color=TEAL_LT)
    add_textbox(
        s, Inches(0.7), Inches(1.65), Inches(12.0), Inches(2.4),
        "Agonistes du GLP-1 :\nmise au point et potentiel\ndans la maladie de Parkinson",
        size=30, bold=True, color=WHITE, font="Georgia",
    )
    add_textbox(
        s, Inches(0.7), Inches(4.25), Inches(12.0), Inches(0.55),
        "Physiologie  ·  Pharmacologie  ·  Indications métaboliques  ·  Repositionnement neurologique",
        size=14, color=RGBColor(0xB0, 0xD0, 0xD8),
    )
    add_textbox(
        s, Inches(0.7), Inches(6.15), Inches(8.2), Inches(0.9),
        "Synthèse bibliographique critique\nPas d’indication actuelle dans Parkinson hors recherche",
        size=13, color=WHITE,
    )
    add_textbox(
        s, Inches(9.2), Inches(6.3), Inches(3.5), Inches(0.7),
        "Présentation de soutenance\n30 diapositives",
        size=12, color=WHITE, align=PP_ALIGN.RIGHT,
    )

    # 2. PLAN
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Organisation", "Plan de la présentation",
           "Du besoin médical aux preuves cliniques — molécule par molécule")
    parts = [
        ("01", "Contexte", "Parkinson, hétérogénéité, impasse modificatrice, logique de repositionnement"),
        ("02", "Chapitre I", "Physiologie du GLP-1, GLP-1R, panorama pharmacocinétique"),
        ("03", "Chapitre II", "DT2, obésité, protection cardio-rénale, tolérance"),
        ("04", "Chapitre III", "Neuroprotection, LIXIPARK, Exenatide-PD3, pharmacien"),
        ("05", "Conclusion", "Pas d’effet de classe ; essais confirmatoires nécessaires"),
    ]
    for i, (n, t, d) in enumerate(parts):
        x = Inches(0.35) + Inches(i * 2.55)
        add_round(s, x, Inches(1.7), Inches(2.4), Inches(4.7), WHITE)
        add_rect(s, x, Inches(1.7), Inches(2.4), Inches(1.15), NAVY if i != 3 else TEAL)
        add_textbox(s, x + Inches(0.15), Inches(1.85), Inches(2.1), Inches(0.4),
                    n, size=28, bold=True, color=ACCENT if i == 3 else TEAL_LT, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.1), Inches(2.35), Inches(2.2), Inches(0.4),
                    t, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.15), Inches(3.15), Inches(2.1), Inches(2.8),
                    d, size=13, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(s, p)

    # 3. OBJECTIFS
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Objectifs", "Objectifs de la thèse",
           "Distinguer plausibilité biologique, résultat préclinique, bénéfice symptomatique et effet modificateur")
    objs = [
        ("1", "Physiologie et pharmacologie",
         "Présenter le GLP-1, le GLP-1R et comparer les agonistes : structure, PK, action courte vs prolongée, tirzépatide non sélectif."),
        ("2", "Indications validées",
         "Analyser les bénéfices dans le DT2, l’obésité et la protection cardio-rénale (LEADER, SUSTAIN-6, REWIND, SELECT, FLOW)."),
        ("3", "Repositionnement Parkinson",
         "Examiner de manière critique préclinique et clinique : LIXIPARK, Exenatide-PD3, MOST-ABLE — sans inférer un effet de classe."),
        ("4", "Rôle du pharmacien",
         "Sécuriser l’usage dans les indications métaboliques et rappeler l’absence d’indication neurologique hors essai."),
    ]
    for i, (n, t, d) in enumerate(objs):
        y = Inches(1.55) + Inches(i * 1.25)
        add_round(s, Inches(0.35), y, Inches(12.6), Inches(1.1), WHITE)
        add_rect(s, Inches(0.35), y, Inches(0.9), Inches(1.1), NAVY if i != 2 else ACCENT)
        add_textbox(s, Inches(0.35), y + Inches(0.28), Inches(0.9), Inches(0.55),
                    n, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.5), y + Inches(0.18), Inches(11), Inches(0.35),
                    t, size=16, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, Inches(1.5), y + Inches(0.55), Inches(11), Inches(0.4),
                    d, size=13, color=MUTED)
    add_footer(s, p)

    # 4. DIVIDER A
    p += 1
    section_divider(
        prs, blank, "A", "Contexte et besoin médical",
        "Maladie de Parkinson : fardeau, hétérogénéité et absence de traitement modificateur", p,
    )

    # 5. PARKINSON — ENJEU + CLINIQUE
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Contexte", "La maladie de Parkinson : un enjeu de santé publique",
           "Deuxième maladie neurodégénérative après Alzheimer — diagnostic clinique et fardeau croissant")
    kpi(s, Inches(0.35), Inches(1.55), Inches(4.0), Inches(1.7), "2ᵉ", "Maladie neurodégénérative", "après Alzheimer")
    kpi(s, Inches(4.55), Inches(1.55), Inches(4.0), Inches(1.7), "Moteur", "Bradykinésie + tremblement ou rigidité", "diagnostic clinique", bg=TEAL)
    kpi(s, Inches(8.75), Inches(1.55), Inches(4.2), Inches(1.7), "Non moteur", "Sommeil, constipation, cognition", "souvent précoces", bg=ACCENT_DK)
    card(s, Inches(0.35), Inches(3.5), Inches(6.2), Inches(3.2),
         "Impact clinique et sociétal",
         ["Autonomie, qualité de vie et charge des aidants",
          "Troubles de la marche, instabilité, complications motrices",
          "Les symptômes non moteurs pèsent autant que le moteur",
          "Fardeau croissant : vieillissement et meilleur diagnostic"],
         TEAL)
    card(s, Inches(6.8), Inches(3.5), Inches(6.15), Inches(3.2),
         "Urgence thérapeutique",
         ["Prise en charge actuelle principalement symptomatique",
          "Aucun traitement n’a démontré un effet modificateur certain",
          "Progression lente, variable, difficile à mesurer",
          "Besoin de biomarqueurs de progression (clinique, imagerie, biologie)"],
         NAVY)
    add_footer(s, p)

    # 6. PHYSIOPATHOLOGIE
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Contexte", "Physiopathologie multifactorielle",
           "La perte dopaminergique n’explique pas à elle seule l’hétérogénéité clinique")
    items = [
        ("Neurones DA / SNpc", "Dégénérescence de la substance noire pars compacta → ↓ dopamine striatale"),
        ("α-synucléine", "Agrégation anormale, toxicité et propagation"),
        ("Mitochondries", "Dysfonction énergétique et vulnérabilité neuronale"),
        ("Stress oxydatif", "Espèces réactives de l’oxygène, lésions cumulatives"),
        ("Neuroinflammation", "Activation microgliale et médiateurs inflammatoires"),
        ("Autophagie / protéostase", "Accumulation de protéines mal conformées"),
    ]
    for i, (t, d) in enumerate(items):
        col, row = i % 3, i // 3
        x = Inches(0.35) + Inches(col * 4.25)
        y = Inches(1.55) + Inches(row * 2.55)
        add_round(s, x, y, Inches(4.05), Inches(2.35), WHITE)
        add_rect(s, x, y, Inches(0.12), Inches(2.35), ACCENT if i in (0, 3) else TEAL)
        add_textbox(s, x + Inches(0.35), y + Inches(0.4), Inches(3.5), Inches(0.55),
                    t, size=16, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, x + Inches(0.35), y + Inches(1.1), Inches(3.5), Inches(0.95),
                    d, size=13, color=MUTED)
    add_footer(s, p)

    # 7. TRAITEMENTS ACTUELS
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Contexte", "Prise en charge actuelle : efficace, mais limitée",
           "La lévodopa reste la référence motrice — sans preuve d’effet modificateur")
    rows = [
        ["Classe", "Rôle principal", "Limite"],
        ["Lévodopa (± ICOMT)", "Référence des symptômes moteurs", "Fluctuations, dyskinésies"],
        ["Agonistes dopaminergiques", "Adaptation au profil du patient", "Effets indésirables neuropsychiatriques"],
        ["IMAO-B", "Potentialisation dopaminergique", "Effet symptomatique seul"],
        ["Amantadine", "Dyskinésies", "Pas d’effet modificateur démontré"],
        ["Stimulation cérébrale profonde", "Complications motrices sélectionnées", "Invasive, patients sélectionnés"],
    ]
    add_table(s, Inches(0.35), Inches(1.55), Inches(12.6), rows,
              [Inches(3.4), Inches(4.6), Inches(4.6)], font_size=13, row_h=0.52)
    add_round(s, Inches(0.35), Inches(5.15), Inches(12.6), Inches(1.55), NAVY)
    add_textbox(s, Inches(0.6), Inches(5.35), Inches(12), Inches(0.35),
                "Message méthodologique", size=13, bold=True, color=ACCENT)
    add_textbox(
        s, Inches(0.6), Inches(5.75), Inches(12), Inches(0.7),
        "Une amélioration du MDS-UPDRS, même à l’état OFF, peut être symptomatique, variable ou indirecte. Elle n’équivaut pas à un ralentissement de la neurodégénérescence.",
        size=14, color=WHITE,
    )
    add_footer(s, p)

    # 8. REPOSITIONNEMENT
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Contexte", "Repositionnement thérapeutique : attractif, non suffisant",
           "Un profil déjà connu ne dispense pas d’une preuve d’efficacité dans la nouvelle indication")
    card(s, Inches(0.35), Inches(1.55), Inches(6.2), Inches(5.15),
         "Intérêt du repositionnement",
         [("Pharmacologie connue", " — PK, toxicologie, pharmacovigilance"),
          ("Développement potentiellement plus court", " — vs nouvelle entité chimique"),
          ("Classe GLP-1 mature", " — DT2, obésité, cardio-rénal"),
          ("Pertinent en neurodégénérescence", " — nombreux échecs de candidats"),
          ("Exigence inchangée", " — efficacité et sécurité dans Parkinson")],
         TEAL, body_size=13)
    card(s, Inches(6.8), Inches(1.55), Inches(6.15), Inches(5.15),
         "Pourquoi le GLP-1 ici ?",
         [("GLP-1R au SNC", " — distribution compatible avec des effets extra-glycémiques"),
          ("Voies de survie", " — AMPc/PKA, PI3K/Akt, MAPK/ERK"),
          ("Mitochondries, oxydatif, inflammation", " — modèles expérimentaux"),
          ("Molécules non équivalentes", " — structure, PK, exposition tissulaire"),
          ("Pas d’effet de classe a priori", " — lire molécule par molécule")],
         NAVY, body_size=13)
    add_footer(s, p)

    # 9. DIVIDER B
    p += 1
    section_divider(
        prs, blank, "B", "Chapitre I — Physiologie et pharmacologie",
        "GLP-1, GLP-1R et différences entre agonistes", p,
    )

    # 10. PHYSIOLOGIE GLP-1
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre I", "Physiologie du GLP-1",
           "Hormone incrétine — biosynthèse, satiété, inactivation en 1–2 minutes")
    add_picture_fit(s, ASSETS / "fig0_physiologie_glp1.png",
                    Inches(0.3), Inches(1.5), Inches(7.15), Inches(4.0))
    add_textbox(s, Inches(0.35), Inches(5.55), Inches(7.1), Inches(0.35),
                "Figure — Production intestinale, effets pancréatiques et satiété centrale.",
                size=11, color=MUTED)
    card(s, Inches(7.6), Inches(1.5), Inches(5.35), Inches(5.2),
         "Points à retenir",
         ["Proglucagon : glucagon (α) vs GLP-1 (cellules L)",
          "Formes actives : GLP-1(7-36)amide et (7-37)",
          "Sécrétion postprandiale ; réponse précoce neuro-hormonale",
          "Insulinosécrétion glucose-dépendante",
          "↓ glucagon si hyperglycémie ; ralentit la vidange",
          "Satiété / ↓ apports — bénéfice ou risque selon le patient",
          "DPP-4 : demi-vie endogène ≈ 1–2 min"],
         TEAL, body_size=12)
    add_footer(s, p)

    # 11. RECEPTEUR + FIGURE 1
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre I", "Récepteur GLP-1R et signalisation",
           "GPCR de famille B — figure 1 de la thèse")
    add_picture_fit(s, ASSETS / "fig1_glp1r_signalisation.png",
                    Inches(0.3), Inches(1.5), Inches(7.35), Inches(4.15))
    add_textbox(s, Inches(0.35), Inches(5.7), Inches(7.3), Inches(0.9),
                "Figure 1 — Gs → AMPc → PKA / EPAC2 (insulinosécrétion). PI3K/Akt et MAPK/ERK : survie et réponse au stress dans les modèles.",
                size=12, color=MUTED)
    add_round(s, Inches(7.85), Inches(1.5), Inches(5.1), Inches(5.2), WHITE)
    add_rect(s, Inches(7.85), Inches(1.5), Inches(5.1), Inches(0.42), NAVY)
    add_textbox(s, Inches(8.0), Inches(1.57), Inches(4.8), Inches(0.3),
                "Distribution et limites", size=13, bold=True, color=WHITE)
    add_bullets(s, Inches(8.05), Inches(2.1), Inches(4.7), Inches(4.4), [
        ("Cellules β", " — expression principale"),
        ("Tube digestif, CV, rein", " — effets extra-glycémiques"),
        ("Certaines régions du SNC", " — compatible, non démonstratif"),
        ("Lecture prudente", " — expression ≠ bénéfice clinique"),
        ("Signalisation expérimentale", " — pas une preuve humaine"),
    ], size=13, spacing=8)
    add_footer(s, p)

    # 12. CLASSIFICATION
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre I", "Classification : pas un groupe homogène",
           "Même récepteur, structures, durées d’action et profils PD différents")
    card(s, Inches(0.35), Inches(1.55), Inches(4.1), Inches(5.15),
         "Origine structurale",
         [("Exendine-4", " — exénatide, lixisénatide"),
          ("Heloderma suspectum", " — homologie partielle au GLP-1"),
          ("Analogues humains", " — liraglutide, dulaglutide, sémaglutide"),
          ("Tirzépatide", " — double agoniste GIP/GLP-1"),
          ("Conséquence", " — ne pas attribuer ses effets au seul GLP-1R")],
         TEAL, body_size=13)
    card(s, Inches(4.65), Inches(1.55), Inches(4.1), Inches(5.15),
         "Action courte",
         [("Exénatide standard, lixisénatide", " — stimulation intermittente"),
          ("Vidange gastrique", " — effet marqué, postprandial"),
          ("Exénatide LP", " — hebdomadaire ; évalué dans PD"),
          ("Lixisénatide", " — 1×/j ; essai LIXIPARK")],
         NAVY, body_size=13)
    card(s, Inches(8.95), Inches(1.55), Inches(4.0), Inches(5.15),
         "Action prolongée",
         [("Lira, dula, séma", " — exposition plus continue"),
          ("Jeûne / HbA1c", " — souvent davantage améliorés"),
          ("Tachyphylaxie gastrique", " — effet vidange peut s’atténuer"),
          ("Pas d’équivalence SNC", " — taille, liaison, distribution")],
         ACCENT_DK, body_size=13)
    add_footer(s, p)

    # 13. TABLEAU 1
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre I", "Tableau 1 — Principaux agonistes et tirzépatide",
           "Indications selon spécialités, doses et cadre réglementaire national")
    rows = [
        ["Molécule", "Structure / mécanisme", "Demi-vie", "Voie", "Particularités"],
        ["Exénatide", "Exendine-4 synthétique, résistant DPP-4", "Quelques heures (std)", "SC 2×/j ; LP 1×/sem.", "LP évalué dans Parkinson"],
        ["Lixisénatide", "Analogue exendine-4, action courte", "~3 h", "SC 1×/j", "Postprandial marqué ; LIXIPARK"],
        ["Liraglutide", "GLP-1 humain + chaîne C16", "~13 h", "SC 1×/j", "DT2 ; obésité selon dose / réglementation"],
        ["Dulaglutide", "2 analogues + Fc IgG4", "~5 j", "SC 1×/sem.", "Bénéfice CV dans REWIND"],
        ["Sémaglutide", "GLP-1 modifié + chaîne C18", "~1 sem.", "SC 1×/sem. ; oral", "DT2, obésité ; données cardio-rénales"],
        ["Tirzépatide", "Double agoniste GIP/GLP-1", "~5 j", "SC 1×/sem.", "Non sélectif du GLP-1R"],
    ]
    add_table(s, Inches(0.2), Inches(1.5), Inches(12.9), rows,
              [Inches(1.7), Inches(3.5), Inches(2.05), Inches(2.35), Inches(3.3)],
              font_size=11, row_h=0.62)
    add_textbox(
        s, Inches(0.35), Inches(6.35), Inches(12.6), Inches(0.45),
        "Point clé : une activité commune sur le GLP-1R ne permet pas de présumer une équivalence des effets neurologiques.",
        size=13, color=MUTED,
    )
    add_footer(s, p)

    # 14. SNC / LIMITES
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre I", "Exposition au SNC et limites de transposition",
           "Plausibilité biologique ≠ démonstration d’efficacité clinique")
    card(s, Inches(0.35), Inches(1.55), Inches(6.2), Inches(5.15),
         "Pourquoi les molécules ne sont pas interchangeables",
         ["Taille, structure, demi-vie variables",
          "Liaison protéique et distribution tissulaire différentes",
          "Accès aux cellules cibles du SNC non superposable",
          "Signal pharmacologique potentiellement distinct",
          "Effets directs (SNC) vs indirects (poids, glycémie, inflammation, CV)"],
         TEAL)
    card(s, Inches(6.8), Inches(1.55), Inches(6.15), Inches(5.15),
         "Dans Parkinson, trois lectures à séparer",
         ["Effet symptomatique (score moteur transitoire)",
          "Effet métabolique ou vasculaire indirect",
          "Véritable effet modificateur de la neurodégénérescence",
          "Les échelles (MDS-UPDRS) restent indispensables",
          "Interprétation : état ON/OFF, variabilité, pertinence clinique"],
         NAVY)
    add_footer(s, p)

    # 15. DIVIDER C
    p += 1
    section_divider(
        prs, blank, "C", "Chapitre II — Applications thérapeutiques",
        "Diabète de type 2, obésité et protection cardio-rénale", p,
    )

    # 16. DT2
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre II", "Diabète de type 2 : place des agonistes du GLP-1R",
           "Action plurimodale — le choix ne se réduit plus à l’HbA1c")
    card(s, Inches(0.35), Inches(1.55), Inches(6.2), Inches(5.15),
         "Rationnel physiopathologique",
         ["Insulinorésistance + altération progressive de l’insulinosécrétion",
          "Dérégulation du glucagon et diminution de l’effet incrétine",
          "Insulinosécrétion glucose-dépendante → faible risque d’hypoglycémie intrinsèque",
          "Risque ↑ si association insuline ou sulfamides",
          "Ralentissement de la vidange et augmentation de la satiété"],
         TEAL)
    card(s, Inches(6.8), Inches(1.55), Inches(6.15), Inches(5.15),
         "Décision individualisée",
         ["Intégrer risque CV, MRC, poids, hypoglycémie, autogestion",
          "Si MCV athéroscléreuse ou haut risque CV : agoniste au bénéfice CV démontré, indépendamment de l’HbA1c initiale ou de la metformine préalable",
          "Association GLP-1R + iSGLT2 possible si risque CV, MRC, IC ou objectif pondéral",
          "Évaluer tolérance, déshydratation, complexité, coût, persistance"],
         NAVY)
    add_footer(s, p)

    # 17. POIDS / OBESITE
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre II", "Contrôle glycémique, poids et obésité",
           "Perte de poids = bénéfice ou risque selon le phénotype")
    card(s, Inches(0.35), Inches(1.55), Inches(4.1), Inches(5.15),
         "Glycémie",
         ["Parmi les non-insuliniques les plus efficaces sur l’HbA1c",
          "Amplitude : glycémie initiale, dose, molécule, adhésion",
          "Action courte → postprandial (vidange)",
          "Action prolongée → jeûne et HbA1c davantage",
          "Améliorations intermédiaires ≠ preuve d’événements CV/rénaux"],
         TEAL)
    card(s, Inches(4.65), Inches(1.55), Inches(4.1), Inches(5.15),
         "Obésité — maladie chronique",
         ["Liraglutide 3,0 mg > placebo + mode de vie",
          "Sémaglutide hebdomadaire : réduction pondérale importante",
          "Tirzépatide : pertes substantielles — mécanisme GIP/GLP-1",
          "Réponse variable ; reprise fréquente à l’arrêt",
          "Objectifs réalistes, réévaluation, conditions d’arrêt"],
         NAVY)
    card(s, Inches(8.95), Inches(1.55), Inches(4.0), Inches(5.15),
         "Vigilance nutritionnelle",
         ["Bénéfice si obésité / DT2 à haut risque cardio-métabolique",
          "Personne âgée, fragile, dénutrie, sarcopénique : risque",
          "Perte de masse maigre, déshydratation, ↓ capacités",
          "Suivre apports, digestif, hydratation, force, autonomie — pas le seul poids"],
         WARN)
    add_footer(s, p)

    # 18. TABLEAU 2 CARDIO-RENAL
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre II", "Tableau 2 — Essais cardio-rénaux",
           "Le bénéfice se formule molécule par molécule, population par population")
    rows = [
        ["Essai", "Molécule", "Population", "Résultat principal"],
        ["LEADER", "Liraglutide", "DT2 à haut risque CV", "Réduction du MACE"],
        ["SUSTAIN-6", "Sémaglutide", "DT2 à haut risque CV", "Résultat favorable sur le MACE"],
        ["REWIND", "Dulaglutide", "DT2, risque CV (dont sans MCV athéroscléreuse documentée)", "Réduction du MACE"],
        ["SELECT", "Sémaglutide", "Surpoids/obésité, MCV établie, sans DT2", "Réduction du MACE"],
        ["FLOW", "Sémaglutide", "DT2 et maladie rénale chronique", "Réduction des événements rénaux majeurs"],
    ]
    add_table(s, Inches(0.25), Inches(1.5), Inches(12.8), rows,
              [Inches(1.7), Inches(2.0), Inches(5.1), Inches(4.0)], font_size=12, row_h=0.62)
    add_round(s, Inches(0.35), Inches(5.45), Inches(12.6), Inches(1.25), SOFT)
    add_textbox(
        s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(0.95),
        "SELECT montre que le bénéfice CV du sémaglutide ne repose pas uniquement sur la glycémie. FLOW conforte une stratégie cardio-rénale individualisée. Pas d’interchangeabilité de classe. La MRC reste multifactorielle (PA, SRAA, iSGLT2, mode de vie).",
        size=13, color=DARK,
    )
    add_footer(s, p)

    # 19. TOLERANCE + PHARMACIEN
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre II", "Tolérance, sécurité et rôle du pharmacien",
           "EI digestifs fréquents à l’initiation — éducation et titration")
    card(s, Inches(0.35), Inches(1.55), Inches(4.1), Inches(5.15),
         "Effets indésirables fréquents",
         ["Nausées, vomissements, diarrhée, constipation",
          "Douleurs abdominales, ↓ appétit",
          "Surtout instauration / hausses de dose",
          "Souvent transitoires, parfois cause d’arrêt",
          "Titration progressive + repas plus petits, moins gras"],
         TEAL)
    card(s, Inches(4.65), Inches(1.55), Inches(4.1), Inches(5.15),
         "Signaux d’alerte",
         ["Douleur abdominale intense persistante",
          "Vomissements sévères, incapacité à s’hydrater",
          "Évoquer complication biliaire ou pancréatite",
          "Hypoglycémie si insuline / sulfamide",
          "Vérifier RCP : CI, précautions, adaptations"],
         WARN)
    card(s, Inches(8.95), Inches(1.55), Inches(4.0), Inches(5.15),
         "Missions du pharmacien",
         ["Technique d’injection, stylo, conservation",
          "Titration, oubli, dextérité, vision, cognition",
          "Adhésion, apports, poids, dénutrition, hydratation",
          "Coordination prescripteur si adaptation associée"],
         NAVY)
    add_footer(s, p)

    # 20. DIVIDER D
    p += 1
    section_divider(
        prs, blank, "D", "Chapitre III — Parkinson et repositionnement",
        "Mécanismes, préclinique, essais cliniques et sécurisation pharmaceutique", p,
    )

    # 21. HETEROGENEITE + MODIFICATEUR
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III", "Besoin non couvert et effet modificateur",
           "Améliorer un score ≠ préserver les neurones")
    card(s, Inches(0.35), Inches(1.55), Inches(6.2), Inches(5.15),
         "Hétérogénéité clinique",
         ["Présentation motrice, non motrice, cognitive, dysautonomie",
          "Vitesse d’évolution et réponse dopaminergique variables",
          "Un modificateur ne devrait pas être jugé au seul score moteur",
          "Besoin double : mieux traiter les symptômes insuffisamment contrôlés et ralentir l’évolution"],
         TEAL)
    card(s, Inches(6.8), Inches(1.55), Inches(6.15), Inches(5.15),
         "Exigences méthodologiques",
         ["Mesures OFF : réduisent l’effet immédiat du traitement, ne prouvent pas la neuroprotection",
          "Durée, insu, traitements associés, critères secondaires",
          "Biomarqueurs exploratoires si possible — aucun biomarqueur unique validé",
          "Efficacité métabolique n’entraîne pas une efficacité neurologique"],
         NAVY)
    add_footer(s, p)

    # 22. MECANISMES / FIGURE 2
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III", "Mécanismes neuroprotecteurs putatifs",
           "Figure 2 — hypothèses précliniques, non une preuve thérapeutique humaine")
    add_picture_fit(s, ASSETS / "fig2_neuroprotection.png",
                    Inches(0.3), Inches(1.48), Inches(8.0), Inches(4.5))
    add_round(s, Inches(8.5), Inches(1.5), Inches(4.45), Inches(5.2), WHITE)
    add_rect(s, Inches(8.5), Inches(1.5), Inches(4.45), Inches(0.42), NAVY)
    add_textbox(s, Inches(8.65), Inches(1.57), Inches(4.15), Inches(0.3),
                "Quatre axes expérimentaux", size=13, bold=True, color=WHITE)
    add_bullets(s, Inches(8.7), Inches(2.1), Inches(4.05), Inches(4.4), [
        ("Survie", " — AMPc/PKA, PI3K/Akt, MAPK/ERK"),
        ("Mitochondries / ROS", " — paramètres améliorés in vitro / in vivo"),
        ("Neuroinflammation", " — ↓ activation microgliale (modèles)"),
        ("Autophagie / α-syn", " — données préliminaires"),
        ("Limite", " — un marqueur ↓ ≠ autonomie ni pronostic"),
    ], size=13, spacing=8)
    add_footer(s, p)

    # 23. PRECLINIQUE
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III", "Données précliniques : preuve de concept, transposition limitée",
           "Modèles toxiques (6-OHDA, MPTP) et modèles α-synucléine")
    rows = [
        ["Molécule", "Observations rapportées selon les protocoles"],
        ["Exénatide", "Effets favorables moteurs, biochimiques ou histologiques"],
        ["Liraglutide", "Signaux protecteurs dans des modèles de toxicité DA"],
        ["Lixisénatide", "Effets favorables dans certains modèles dopaminergiques"],
        ["Sémaglutide", "Signaux sur survie, autophagie ou α-synucléine selon les travaux"],
    ]
    add_table(s, Inches(0.35), Inches(1.5), Inches(12.6), rows,
              [Inches(2.5), Inches(10.1)], font_size=13, row_h=0.58)
    add_round(s, Inches(0.35), Inches(4.7), Inches(12.6), Inches(2.0), NAVY)
    add_textbox(s, Inches(0.6), Inches(4.85), Inches(12.1), Inches(0.3),
                "Limites de transposition", size=13, bold=True, color=ACCENT)
    add_textbox(
        s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(1.25),
        "Lésion souvent rapide ; reproduction partielle de la lenteur, de l’hétérogénéité, des non-moteurs et de la pathologie α-syn humaine. Une meilleure performance motrice animale peut être fonctionnelle ou symptomatique. Ces données justifient des essais rigoureux — elles ne valident pas l’efficacité.",
        size=13, color=WHITE,
    )
    add_footer(s, p)

    # 24. LIXIPARK
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III — Clinique", "LIXIPARK — lixisénatide, phase II",
           "Signal moteur favorable à court terme — pas une preuve de neuroprotection durable")
    kpi(s, Inches(0.35), Inches(1.55), Inches(4.0), Inches(1.85), "156", "Patients", "Parkinson précoce", bg=TEAL)
    kpi(s, Inches(4.55), Inches(1.55), Inches(4.0), Inches(1.85), "12 mois", "Durée", "critère moteur", bg=NAVY)
    kpi(s, Inches(8.75), Inches(1.55), Inches(4.2), Inches(1.85), "3,08 pts", "Différence ajustée", "MDS-UPDRS III", bg=ACCENT_DK)
    add_round(s, Inches(0.35), Inches(3.6), Inches(12.6), Inches(3.1), WHITE)
    add_rect(s, Inches(0.35), Inches(3.6), Inches(12.6), Inches(0.45), NAVY)
    add_textbox(s, Inches(0.55), Inches(3.68), Inches(12.2), Inches(0.3),
                "Résultat et lecture critique", size=14, bold=True, color=WHITE)
    add_bullets(s, Inches(0.55), Inches(4.2), Inches(12.2), Inches(2.3), [
        "Variation moyenne MDS-UPDRS III : −0,04 (lixisénatide) vs +3,04 (placebo)",
        "Différence ajustée : 3,08 points en faveur du lixisénatide",
        "Nausées et vomissements plus fréquents sous lixisénatide",
        "Phase II, durée limitée, critère principalement moteur → ne permet pas d’affirmer un ralentissement durable de la neurodégénérescence",
        "Pertinence clinique à discuter : variabilité des scores, évolution naturelle, absence de biomarqueur validé de progression",
    ], size=13, spacing=4)
    add_footer(s, p)

    # 25. EXENATIDE-PD3
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III — Clinique", "Exenatide-PD3 — phase III négative",
           "Donnée majeure : renoncer à toute généralisation optimiste de classe")
    kpi(s, Inches(0.35), Inches(1.55), Inches(4.0), Inches(1.85), "194", "Participants", "exénatide hebdomadaire", bg=NAVY)
    kpi(s, Inches(4.55), Inches(1.55), Inches(4.0), Inches(1.85), "96 sem.", "Suivi", "vs placebo", bg=TEAL)
    kpi(s, Inches(8.75), Inches(1.55), Inches(4.2), Inches(1.85), "Négatif", "Critère principal", "MDS-UPDRS III OFF", bg=WARN)
    card(s, Inches(0.35), Inches(3.65), Inches(6.2), Inches(3.05),
         "Ce qui a précédé",
         ["Étude pilote : signal moteur exploratoire, faible effectif",
          "Essai randomisé de phase II : signal MDS-UPDRS III OFF",
          "Puissance et durée insuffisantes pour un effet neuroprotecteur durable",
          "Ces signaux ont justifié une phase III"],
         TEAL)
    card(s, Inches(6.8), Inches(3.65), Inches(6.15), Inches(3.05),
         "Résultat de la phase III",
         ["Aggravation motrice similaire dans les deux groupes",
          "Aucune différence significative sur le critère principal",
          "Pas d’argument en faveur d’un effet modificateur de l’exénatide dans la population étudiée",
          "N’autorise pas d’extrapoler un échec (ni un succès) à toute la classe"],
         WARN)
    add_footer(s, p)

    # 26. TABLEAU 3 + FIGURE ESSAIS
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III — Clinique", "Tableau 3 — Synthèse des essais dans Parkinson",
           "Résultats hétérogènes — pas d’effet de classe démontré")
    add_picture_fit(s, ASSETS / "fig3_essais_cliniques.png",
                    Inches(0.25), Inches(1.48), Inches(6.35), Inches(3.55))
    rows = [
        ["Étude", "Population / durée", "Résultat", "Limites"],
        ["Exénatide, pilote", "Faible effectif", "Signal exploratoire", "Symptomatique non exclu"],
        ["Exénatide, ph. II", "Hebdomadaire", "Signal OFF", "Puissance limitée"],
        ["LIXIPARK", "156 ; précoce ; 12 mois", "+3,08 pts MDS-UPDRS III", "Ph. II ; EI digestifs"],
        ["Exenatide-PD3", "194 ; 96 semaines", "Pas de différence OFF", "Négatif pour l’exénatide"],
    ]
    add_table(s, Inches(6.7), Inches(1.5), Inches(6.25), rows,
              [Inches(1.55), Inches(1.7), Inches(1.55), Inches(1.45)],
              font_size=10, row_h=0.72)
    add_round(s, Inches(0.35), Inches(5.2), Inches(12.6), Inches(1.5), SOFT)
    add_textbox(
        s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(1.1),
        "Méta-analyse récente d’essais contre placebo : pas de bénéfice statistiquement significatif sur l’ensemble des symptômes moteurs. Un signal exploratoire pour les agents à action courte, fondé sur peu d’études, ne modifie pas les pratiques. Non moteurs et qualité de vie : pas d’amélioration démontrée.",
        size=13, color=DARK,
    )
    add_footer(s, p)

    # 27. PERSPECTIVES / MOST-ABLE
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III — Perspectives", "Que doivent faire les essais futurs ?",
           "MOST-ABLE illustre la poursuite de la recherche — sans résultat d’efficacité à ce jour")
    pers = [
        ("Populations", "Plus homogènes, idéalement stade précoce ; durée compatible avec une progression lente"),
        ("Critères", "OFF moteur + non moteurs, QdV, cognition, autonomie, chutes, nutrition, biomarqueurs exploratoires"),
        ("Symptomatique vs modificateur", "Sevrage, mesures répétées, standardisation OFF, imagerie/biologie — aucun biomarqueur unique validé"),
        ("Stratification", "Phénotypes moteurs, cognitifs, métaboliques ou génétiques — à tester prospectivement, pas a posteriori"),
        ("MOST-ABLE", "Ph. II Japon, 99 patients, PD précoce, 2 doses de sémaglutide oral, 48 semaines, MDS-UPDRS III OFF + QdV, cognition, imagerie DA"),
        ("Message", "Différences PK ou de stade peuvent expliquer l’hétérogénéité — elles ne maintiennent pas une hypothèse sans preuve"),
    ]
    for i, (t, d) in enumerate(pers):
        col, row = i % 3, i // 3
        x = Inches(0.35) + Inches(col * 4.25)
        y = Inches(1.55) + Inches(row * 2.55)
        add_round(s, x, y, Inches(4.05), Inches(2.35), WHITE)
        add_rect(s, x, y, Inches(4.05), Inches(0.5), TEAL if row == 0 else NAVY)
        add_textbox(s, x + Inches(0.15), y + Inches(0.1), Inches(3.75), Inches(0.35),
                    t, size=13, bold=True, color=WHITE)
        add_textbox(s, x + Inches(0.15), y + Inches(0.7), Inches(3.75), Inches(1.45),
                    d, size=12, color=DARK)
    add_footer(s, p)

    # 28. PHARMACIEN PARKINSON
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Chapitre III — Pratique", "Tolérance dans Parkinson et rôle du pharmacien",
           "Patients souvent âgés, polymédiqués, vulnérables sur le plan nutritionnel et fonctionnel")
    roles = [
        ("Risques spécifiques PD",
         ["Constipation déjà fréquente (dysautonomie)",
          "Nausées, vomissements, ↓ appétit",
          "Perte de poids, sarcopénie, hypotension, chutes",
          "Sélection future : éviter dénutrition, dysphagie, constipation sévère"]),
        ("Indications validées",
         ["Injection, conservation, titration, oubli",
          "Repérer EI digestifs et déshydratation",
          "Hypoglycémie si insuline / sulfamides",
          "Tremblements, dextérité, cognition, aidant"]),
        ("Message de santé publique",
         ["Pas d’indication Parkinson hors essai clinique",
          "Prévenir automédication et hors AMM injustifié",
          "Attentes disproportionnées (préclinique, médias)",
          "Coordination neuro / MG / diabéto / diététicien"]),
    ]
    for i, (t, bullets) in enumerate(roles):
        x = Inches(0.35) + Inches(i * 4.3)
        card(s, x, Inches(1.55), Inches(4.1), Inches(5.15), t, bullets,
             TEAL if i == 0 else (WARN if i == 2 else NAVY), body_size=13)
    add_footer(s, p)

    # 29. CONCLUSION
    p += 1
    s = prs.slides.add_slide(blank)
    header(s, "Conclusion générale", "Messages à retenir",
           "Piste biologiquement plausible — encore incomplètement validée")
    msgs = [
        ("1", "Classe métabolique mature",
         "Incontournable dans le DT2 ; obésité et cardio-rénal pour certaines molécules — jamais parfaitement interchangeables."),
        ("2", "Rationnel neurologique",
         "Survie cellulaire, mitochondries, oxydatif, inflammation, autophagie : hypothèses précliniques, pas une preuve humaine."),
        ("3", "Clinique hétérogène",
         "LIXIPARK : signal moteur à 12 mois. Exenatide-PD3 : pas de bénéfice sur la progression. Pas d’effet de classe."),
        ("4", "Pas d’indication Parkinson",
         "Hors protocole de recherche. La place éventuelle dépendra d’essais plus longs, plus puissants, mieux caractérisés."),
        ("5", "Pharmacien",
         "Éducation, vigilance digestive et nutritionnelle, coordination — et information claire sur le caractère expérimental."),
    ]
    for i, (n, t, d) in enumerate(msgs):
        y = Inches(1.48) + Inches(i * 1.05)
        add_round(s, Inches(0.35), y, Inches(12.6), Inches(0.95), WHITE)
        add_rect(s, Inches(0.35), y, Inches(0.7), Inches(0.95), ACCENT if i == 3 else NAVY)
        add_textbox(s, Inches(0.35), y + Inches(0.22), Inches(0.7), Inches(0.5),
                    n, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, Inches(1.25), y + Inches(0.1), Inches(11.4), Inches(0.32),
                    t, size=14, bold=True, color=NAVY, font="Georgia")
        add_textbox(s, Inches(1.25), y + Inches(0.45), Inches(11.4), Inches(0.4),
                    d, size=12, color=MUTED)
    add_footer(s, p)

    # 30. MERCI
    p += 1
    s = prs.slides.add_slide(blank)
    add_bg(s, NAVY)
    add_rect(s, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
    add_textbox(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.4),
                "DISCUSSION", size=14, bold=True, color=ACCENT)
    add_textbox(s, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.4),
                "Merci de votre attention",
                size=40, bold=True, color=WHITE, font="Georgia")
    add_textbox(
        s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.2),
        "Mots-clés : agonistes du GLP-1  ·  maladie de Parkinson  ·  neuroprotection\n"
        "repositionnement thérapeutique  ·  diabète de type 2",
        size=16, color=RGBColor(0xB0, 0xD0, 0xD8),
    )
    add_textbox(
        s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.9),
        "À ce stade, les agonistes du GLP-1R sont des candidats au repositionnement,\n"
        "et non des traitements établis de la maladie de Parkinson.",
        size=16, color=WHITE,
    )
    add_textbox(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.3),
                f"{p}  /  {TOTAL}", size=11, color=RGBColor(0x70, 0x98, 0xA0))

    assert p == TOTAL, f"Expected {TOTAL} slides, got {p}"

    out = Path("/workspace/docs/presentation/Soutenance_GLP1_Parkinson.pptx")
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)} (counter={p})")
    return str(out)


if __name__ == "__main__":
    build()
